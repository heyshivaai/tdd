"""
Document reader: extracts text from PDF, DOCX, XLSX, PPTX, and plain-text
files, then chunks them for Claude API calls.

Why: VDR documents arrive in mixed formats. We normalise everything to text
chunks with source metadata so signals can be traced back to their document.
Chunks stay within Claude's context window (max_chars=32000 per chunk).

Architecture:
  route_document(filepath) → raw extraction dict (text + metadata)
  extract_text(filepath)   → chunked output (backward-compatible entry point)

route_document() is the master entry point for all document extraction.
It detects file type and routes to the correct extractor. Every extractor
returns the same dict shape so callers never need conditional logic.

Skip tracking: every file that cannot be read produces a structured skip
record (filename, reason, file_type) so practitioners see exactly what was
missed and why — no silent failures.
"""
import logging
import mimetypes
import os
from pathlib import Path
from typing import Dict, List, Optional

from pdfminer.high_level import extract_text as pdfminer_extract_text

logger = logging.getLogger(__name__)

MAX_CHARS = 32000

# ── File type classification ─────────────────────────────────────────────────
# Non-extractable: we log them and move on — no crash, no silent skip.
NON_EXTRACTABLE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".svg",
    ".mp4", ".mp3", ".wav", ".m4a", ".avi", ".mov",
    ".zip", ".rar", ".tar", ".gz",
    ".exe", ".dll", ".bin",
}

# Audio/video subset — stub for future Whisper integration; check for sidecar
# transcripts (.txt / .vtt beside the media file) before giving up.
AUDIO_VIDEO_EXTENSIONS = {".mp3", ".mp4", ".wav", ".m4a", ".avi", ".mov"}

# ── Skip reason constants ─────────────────────────────────────────────────────
SKIP_FILE_NOT_FOUND = "File not found on disk"
SKIP_UNSUPPORTED_TYPE = "Unsupported file type"
SKIP_MISSING_DEPENDENCY = "Required library not installed"
SKIP_READ_ERROR = "Failed to read file"
SKIP_NO_TEXT = "File readable but contained no extractable text"
SKIP_ENCRYPTED = "File is encrypted or password-protected"


def make_skip_record(
    filepath: str,
    reason: str,
    detail: str = "",
) -> Dict[str, str]:
    """
    Build a structured skip record for a document that could not be read.

    Args:
        filepath: Path to the file that was skipped.
        reason: One of the SKIP_* constants above.
        detail: Optional extra context (e.g. the exception message).

    Returns:
        Dict with keys: filename, filepath, reason, detail, file_type.
    """
    path = Path(filepath)
    return {
        "filename": path.name,
        "filepath": str(filepath),
        "reason": reason,
        "detail": detail,
        "file_type": path.suffix.lower().lstrip("."),
    }


# Module-level collector so callers can retrieve skipped docs after extraction.
# Reset via reset_skip_log() before each scan run.
_skip_log: List[Dict[str, str]] = []


def reset_skip_log() -> None:
    """Clear the skip log. Call once at the start of each scan run."""
    _skip_log.clear()


def get_skip_log() -> List[Dict[str, str]]:
    """Return a copy of all skip records accumulated since last reset."""
    return list(_skip_log)


def _record_skip(filepath: str, reason: str, detail: str = "") -> None:
    """Log a warning AND append a structured skip record."""
    rec = make_skip_record(filepath, reason, detail)
    _skip_log.append(rec)
    logger.warning("SKIPPED %s — %s%s", rec["filename"], reason, f": {detail}" if detail else "")


def probe_readable(filepath: str) -> bool:
    """
    Lightweight readability check — can we open this file and get *any* text?

    Unlike extract_text(), this reads only the first page/sheet/1KB to verify
    the file is readable, then stops. Used by the VDR preview to surface
    unreadable files without the cost of full extraction.

    Records a skip if the file cannot be read (same skip log as extract_text).

    Args:
        filepath: Path to the file to probe.

    Returns:
        True if the file is readable and contains text, False otherwise.
    """
    path = Path(filepath)
    if not path.exists():
        _record_skip(filepath, SKIP_FILE_NOT_FOUND)
        return False

    ext = path.suffix.lower()

    try:
        if ext == ".pdf":
            return _probe_pdf(filepath)
        elif ext == ".docx":
            return _probe_docx(filepath)
        elif ext in (".xlsx", ".xls"):
            return _probe_excel(filepath)
        elif ext in (".txt", ".md", ".csv"):
            return _probe_plaintext(filepath)
        else:
            _record_skip(filepath, SKIP_UNSUPPORTED_TYPE, f"extension: {ext}")
            return False
    except Exception as exc:
        _record_skip(filepath, SKIP_READ_ERROR, str(exc))
        return False


def _probe_pdf(filepath: str) -> bool:
    """
    Read only the first page of a PDF to check readability using pdfminer.six.

    Extracts text from the first page only (page_numbers=[0]) to keep the probe
    lightweight. If pdfminer returns any text, we consider the file readable.
    """
    try:
        text = pdfminer_extract_text(filepath, page_numbers=[0]).strip()
        if not text:
            _record_skip(filepath, SKIP_NO_TEXT, "Possibly a scanned-image PDF without OCR")
            return False
        return True
    except Exception as exc:
        err_msg = str(exc).lower()
        if "password" in err_msg or "encrypt" in err_msg:
            _record_skip(filepath, SKIP_ENCRYPTED, "PDF requires a password")
        else:
            _record_skip(filepath, SKIP_READ_ERROR, str(exc))
        return False


def _probe_docx(filepath: str) -> bool:
    """Open a DOCX and check if it has any paragraph text."""
    try:
        from docx import Document
    except ImportError:
        _record_skip(filepath, SKIP_MISSING_DEPENDENCY, "pip install python-docx")
        return False
    try:
        doc = Document(filepath)
        for para in doc.paragraphs[:10]:
            if para.text.strip():
                return True
        _record_skip(filepath, SKIP_NO_TEXT)
        return False
    except Exception as exc:
        _record_skip(filepath, SKIP_READ_ERROR, str(exc))
        return False


def _probe_excel(filepath: str) -> bool:
    """Open an Excel file and check if the first sheet has any data."""
    ext = Path(filepath).suffix.lower()
    if ext == ".xls":
        try:
            import xlrd
        except ImportError:
            _record_skip(filepath, SKIP_MISSING_DEPENDENCY, "pip install xlrd")
            return False
        try:
            wb = xlrd.open_workbook(filepath)
            sheet = wb.sheet_by_index(0)
            return sheet.nrows > 0
        except Exception as exc:
            _record_skip(filepath, SKIP_READ_ERROR, str(exc))
            return False
    else:
        try:
            import openpyxl
        except ImportError:
            _record_skip(filepath, SKIP_MISSING_DEPENDENCY, "pip install openpyxl")
            return False
        try:
            wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
            ws = wb[wb.sheetnames[0]]
            for row in ws.iter_rows(max_row=5, values_only=True):
                if any(c is not None for c in row):
                    wb.close()
                    return True
            wb.close()
            _record_skip(filepath, SKIP_NO_TEXT)
            return False
        except Exception as exc:
            _record_skip(filepath, SKIP_READ_ERROR, str(exc))
            return False


def _probe_plaintext(filepath: str) -> bool:
    """Read first 1KB of a text file to check readability."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            sample = f.read(1024).strip()
        if not sample:
            _record_skip(filepath, SKIP_NO_TEXT)
            return False
        return True
    except Exception as exc:
        _record_skip(filepath, SKIP_READ_ERROR, str(exc))
        return False


def route_document(file_path: str) -> dict:
    """
    Master entry point for all document extraction.

    Detects file type and routes to the correct extractor. Always returns the
    same dict shape — callers never need to know which extractor was used.

    Return shape:
    {
        "text": str,                    # extracted text (empty string if not extractable)
        "extraction_method": str,       # "pdfminer" | "python-docx" | "openpyxl" | "python-pptx" | ...
        "quality": str,                 # "good" | "low" | "empty" | "error" | "non_extractable"
        "word_count": int,
        "file_path": str,
        "extraction_error": str | None,
        "metadata": dict                # chunk metadata for signal extraction prompts
    }

    Args:
        file_path: Path to the document to extract text from.

    Returns:
        Dict with the extraction result (always has the same keys).
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    # Base result shape — callers can always rely on these keys existing
    base: Dict = {
        "text": "",
        "extraction_method": "none",
        "quality": "non_extractable",
        "word_count": 0,
        "file_path": file_path,
        "extraction_error": None,
        "metadata": {
            "filename": path.name,
            "extension": ext,
            "file_size_kb": round(path.stat().st_size / 1024, 1) if path.exists() else 0,
        },
    }

    if not path.exists():
        base["quality"] = "error"
        base["extraction_error"] = "File not found"
        _record_skip(file_path, SKIP_FILE_NOT_FOUND)
        return base

    # Non-extractable: images, archives (but not audio/video — those get sidecar check)
    if ext in NON_EXTRACTABLE_EXTENSIONS and ext not in AUDIO_VIDEO_EXTENSIONS:
        base["quality"] = "non_extractable"
        return base

    # Audio/video stub — check for sidecar transcript first
    if ext in AUDIO_VIDEO_EXTENSIONS:
        transcript_candidates = [
            path.with_suffix(".txt"),
            path.with_suffix(".vtt"),
            path.with_name(path.stem + "_transcript.txt"),
        ]
        for candidate in transcript_candidates:
            if candidate.exists():
                try:
                    text = candidate.read_text(encoding="utf-8", errors="replace").strip()
                    base["text"] = text
                    base["extraction_method"] = "transcript_sidecar"
                    base["quality"] = "good" if len(text.split()) > 50 else "low"
                    base["word_count"] = len(text.split())
                    return base
                except Exception as e:
                    base["extraction_error"] = str(e)
        # No transcript found — flag for manual processing
        base["extraction_method"] = "audio_stub"
        base["quality"] = "non_extractable"
        base["metadata"]["note"] = "Audio file with no transcript. Add .txt sidecar or integrate Whisper."
        return base

    # PDF
    if ext == ".pdf":
        result = extract_pdf_text(file_path)
        base.update(result)
        return base

    # DOCX
    if ext in (".docx", ".doc"):
        result = extract_docx_text(file_path)
        base.update(result)
        return base

    # Excel / CSV
    if ext in (".xlsx", ".xls", ".xlsm", ".csv"):
        result = extract_tabular_text(file_path)
        base.update(result)
        return base

    # PowerPoint
    if ext in (".pptx", ".ppt"):
        result = extract_pptx_text(file_path)
        base.update(result)
        return base

    # Plain text fallback
    if ext in (".txt", ".md", ".log", ".json", ".xml", ".html"):
        try:
            text = Path(file_path).read_text(encoding="utf-8", errors="replace").strip()
            base["text"] = text
            base["extraction_method"] = "plaintext"
            base["quality"] = "good" if text else "empty"
            base["word_count"] = len(text.split())
            return base
        except Exception as e:
            base["extraction_error"] = str(e)
            base["quality"] = "error"
            return base

    # Unknown extension — attempt plaintext, flag it
    try:
        text = Path(file_path).read_text(encoding="utf-8", errors="replace").strip()
        base["text"] = text
        base["extraction_method"] = "plaintext_fallback"
        base["quality"] = "low"
        base["word_count"] = len(text.split())
        base["metadata"]["note"] = f"Unknown extension {ext} — read as plaintext"
    except Exception as e:
        base["extraction_error"] = str(e)
        base["quality"] = "error"

    return base


# ── Raw extractors (return dicts, NOT chunked) ──────────────────────────────
# These are called by route_document(). For backward compat, the old
# extract_text() function below still works — it calls route_document()
# internally and then chunks the result.


def extract_pdf_text(file_path: str) -> dict:
    """
    Extract text from a PDF using pdfminer.six.

    Returns a dict with text, extraction_method, quality, and word_count.
    Falls back gracefully if extraction fails.
    """
    result: Dict = {
        "text": "",
        "extraction_method": "pdfminer",
        "quality": "unknown",
        "word_count": 0,
        "file_path": file_path,
        "extraction_error": None,
    }
    try:
        raw_text = pdfminer_extract_text(file_path)
        if raw_text:
            cleaned = raw_text.strip()
            word_count = len(cleaned.split())
            file_size_kb = os.path.getsize(file_path) / 1024
            result["text"] = cleaned
            result["word_count"] = word_count
            # Quality check: a text-heavy PDF should yield >= 50 words per 100KB
            expected_min_words = max(50, (file_size_kb / 100) * 50)
            if word_count == 0:
                result["quality"] = "empty"
            elif word_count < expected_min_words and file_size_kb > 100:
                result["quality"] = "low"  # Likely scanned or image-heavy
                logger.warning(
                    "Low text yield for %s: %d words from %.0f KB file — "
                    "may be scanned/image-heavy",
                    Path(file_path).name, word_count, file_size_kb,
                )
            else:
                result["quality"] = "good"
        else:
            result["quality"] = "empty"
            result["text"] = ""
    except Exception as e:
        err_msg = str(e).lower()
        if "password" in err_msg or "encrypt" in err_msg:
            result["extraction_error"] = "PDF requires a password"
        else:
            result["extraction_error"] = str(e)
        result["quality"] = "error"
        result["text"] = ""
    return result


def extract_docx_text(file_path: str) -> dict:
    """
    Extract text from DOCX preserving section structure.

    Pulls text from paragraphs and tables. Table rows are pipe-delimited
    so Claude can interpret column context.
    """
    result: Dict = {
        "text": "",
        "extraction_method": "python-docx",
        "quality": "unknown",
        "word_count": 0,
        "file_path": file_path,
        "extraction_error": None,
    }
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        # Also pull text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)

        text = "\n\n".join(paragraphs)
        result["text"] = text
        result["word_count"] = len(text.split())
        result["quality"] = "good" if result["word_count"] > 20 else "empty"
    except ImportError:
        result["extraction_error"] = "python-docx not installed"
        result["quality"] = "error"
    except Exception as e:
        result["extraction_error"] = str(e)
        result["quality"] = "error"
    return result


def extract_tabular_text(file_path: str, max_rows: int = 0) -> dict:
    """
    Extract text from Excel/CSV preserving headers and row structure.

    Does NOT flatten to a string blob — preserves tabular shape so
    Claude can interpret column context correctly. Each sheet gets its
    own section header.

    Args:
        file_path: Path to the Excel or CSV file.
        max_rows: Maximum rows to read per sheet. 0 = no limit (read all rows).
                  For very large files, the 500K char safety truncation still
                  applies to prevent downstream memory issues.
    """
    result: Dict = {
        "text": "",
        "extraction_method": "openpyxl",
        "quality": "unknown",
        "word_count": 0,
        "file_path": file_path,
        "extraction_error": None,
    }
    ext = Path(file_path).suffix.lower()
    nrows_arg = max_rows if max_rows > 0 else None  # None = read all rows
    try:
        import pandas as pd

        if ext == ".csv":
            df_dict = {"Sheet1": pd.read_csv(file_path, nrows=nrows_arg, dtype=str).fillna("")}
        elif ext == ".xls":
            df_dict = pd.read_excel(
                file_path, sheet_name=None, nrows=nrows_arg, engine="xlrd", dtype=str
            )
            df_dict = {k: v.fillna("") for k, v in df_dict.items()}
        else:
            df_dict = pd.read_excel(
                file_path, sheet_name=None, nrows=nrows_arg, engine="openpyxl", dtype=str
            )
            df_dict = {k: v.fillna("") for k, v in df_dict.items()}

        sections = []
        for sheet_name, df in df_dict.items():
            if df.empty:
                continue
            header_line = " | ".join(str(c) for c in df.columns)
            rows = []
            for _, row in df.iterrows():
                row_text = " | ".join(str(v) for v in row.values if str(v).strip())
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sections.append(
                    f"[Sheet: {sheet_name}]\nHeaders: {header_line}\n" + "\n".join(rows)
                )

        full_text = "\n\n".join(sections)

        # Safety: truncate if extracted text exceeds 500K chars to prevent
        # downstream memory issues (Claude context window limit).
        if len(full_text) > 500_000:
            logger.warning(
                "Excel/CSV file %s extracted text exceeds 500K chars (%d chars) — truncating",
                Path(file_path).name, len(full_text),
            )
            full_text = full_text[:500_000] + "\n[...truncated due to size]"

        result["text"] = full_text
        result["word_count"] = len(result["text"].split())
        result["quality"] = "good" if result["word_count"] > 10 else "empty"
        result["extraction_method"] = "pandas+openpyxl"
    except ImportError as e:
        result["extraction_error"] = f"Missing dependency: {e}"
        result["quality"] = "error"
    except Exception as e:
        result["extraction_error"] = str(e)
        result["quality"] = "error"
    return result


def extract_pptx_text(file_path: str) -> dict:
    """
    Extract text from PowerPoint slides.

    Each slide is labelled [Slide N] so Claude knows where in the deck
    a signal came from. Pulls text from all shapes with text frames.
    """
    result: Dict = {
        "text": "",
        "extraction_method": "python-pptx",
        "quality": "unknown",
        "word_count": 0,
        "file_path": file_path,
        "extraction_error": None,
    }
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
            if slide_parts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_parts))

        result["text"] = "\n\n".join(slides_text)
        result["word_count"] = len(result["text"].split())
        result["quality"] = "good" if result["word_count"] > 20 else "empty"
    except ImportError:
        result["extraction_error"] = "python-pptx not installed"
        result["quality"] = "error"
    except Exception as e:
        result["extraction_error"] = str(e)
        result["quality"] = "error"
    return result


# ── Backward-compatible entry points ────────────────────────────────────────
# These wrap route_document() + chunk_text() so existing callers (vdr_triage,
# domain_analyst, resume_horizon) keep working without changes.


def extract_text(filepath: str) -> List[dict]:
    """
    Extract text from any supported file type and return as a list of chunks.

    This is the backward-compatible entry point. It calls route_document()
    to get raw text, then chunks it via chunk_text(). Callers that need the
    raw extraction dict (with quality, metadata, etc.) should call
    route_document() directly instead.

    Files that cannot be read are recorded in the module skip log (see
    get_skip_log()) with a structured reason so the pipeline can surface
    them to practitioners.

    Args:
        filepath: Path to the file to extract text from.

    Returns:
        List of chunk dicts, or empty list if file cannot be read.
    """
    extracted = route_document(filepath)

    # Record skips for files that failed or had no text
    if extracted["quality"] == "error" and extracted["extraction_error"]:
        err = extracted["extraction_error"]
        if "password" in err.lower() or "encrypt" in err.lower():
            _record_skip(filepath, SKIP_ENCRYPTED, err)
        elif "not installed" in err.lower() or "missing dependency" in err.lower():
            _record_skip(filepath, SKIP_MISSING_DEPENDENCY, err)
        else:
            _record_skip(filepath, SKIP_READ_ERROR, err)
        return []

    if extracted["quality"] == "non_extractable":
        ext = Path(filepath).suffix.lower()
        _record_skip(filepath, SKIP_UNSUPPORTED_TYPE, f"extension: {ext}")
        return []

    if extracted["quality"] == "empty" or not extracted["text"]:
        _record_skip(filepath, SKIP_NO_TEXT)
        return []

    return chunk_text(extracted["text"], source_doc=Path(filepath).name)


def extract_text_from_pdf(filepath: str) -> List[dict]:
    """
    Extract text from a PDF and return as a list of chunks.

    Backward-compatible wrapper around extract_pdf_text() + chunk_text().
    Each chunk is a dict with keys: text, source_doc, chunk_index, total_chunks.
    Returns an empty list if the file cannot be read (logs the error).

    Args:
        filepath: Path to the PDF file to extract text from.

    Returns:
        List of chunk dicts, or empty list if file cannot be read.
    """
    path = Path(filepath)
    if not path.exists():
        _record_skip(filepath, SKIP_FILE_NOT_FOUND)
        return []

    raw = extract_pdf_text(filepath)

    if raw["quality"] == "error":
        err = raw["extraction_error"] or "Unknown PDF read error"
        if "password" in err.lower():
            _record_skip(filepath, SKIP_ENCRYPTED, err)
        else:
            _record_skip(filepath, SKIP_READ_ERROR, err)
        return []

    if not raw["text"]:
        _record_skip(filepath, SKIP_NO_TEXT, "Possibly a scanned-image PDF without OCR")
        return []

    return chunk_text(raw["text"], source_doc=path.name)


def chunk_text(text: str, source_doc: str, max_chars: int = MAX_CHARS) -> List[dict]:
    """
    Split text into chunks no larger than max_chars, preserving word boundaries.

    This is the legacy chunker kept for backward compatibility. New code should
    prefer chunk_document() which is structure-aware (respects page, slide,
    sheet, and paragraph boundaries).

    Never breaks a word. If a single word exceeds max_chars, logs a warning but
    includes it in its own chunk (to ensure no data is lost).

    Args:
        text: The text to chunk.
        source_doc: The source document name (for metadata).
        max_chars: Maximum character limit per chunk (default 32000).

    Returns:
        List of dicts: {text, source_doc, chunk_index, total_chunks}.
    """
    words = text.split()
    chunks: List[str] = []
    current_chunk_words: List[str] = []
    current_length = 0

    for word in words:
        word_len = len(word) + 1  # +1 for space
        if current_length + word_len > max_chars and current_chunk_words:
            chunks.append(" ".join(current_chunk_words))
            current_chunk_words = [word]
            current_length = word_len
        else:
            current_chunk_words.append(word)
            current_length += word_len

    if current_chunk_words:
        chunks.append(" ".join(current_chunk_words))

    total = len(chunks)
    return [
        {
            "text": chunk,
            "source_doc": source_doc,
            "chunk_index": i,
            "total_chunks": total,
        }
        for i, chunk in enumerate(chunks)
    ]


# ── Structure-aware chunking ────────────────────────────────────────────────
# These functions work with route_document() output and produce richer chunk
# metadata (section_hint, extraction_method, quality) for signal extraction.


def chunk_document(extracted: dict, max_chars: int = MAX_CHARS) -> List[dict]:
    """
    Split extracted document text into chunks, respecting document structure.

    Unlike chunk_text() which splits on word boundaries only, this function:
      - PDFs: chunks at paragraph (double-newline) boundaries
      - DOCX: chunks at paragraph boundaries
      - Excel: chunks at sheet boundaries ([Sheet: ...] markers)
      - PPTX: chunks at slide boundaries ([Slide N] markers)
      - Everything else: paragraph-aware split

    Each chunk carries metadata so Claude knows its position in the doc.

    Args:
        extracted: Dict from route_document() with text, extraction_method, etc.
        max_chars: Maximum characters per chunk (default 32000).

    Returns:
        List of chunk dicts with keys: text, chunk_index, total_chunks,
        source_doc, section_hint, extraction_method, quality.
        Empty list if extracted text is empty.
    """
    text = extracted.get("text", "")
    source_doc = extracted.get("file_path", "unknown")
    method = extracted.get("extraction_method", "unknown")
    quality = extracted.get("quality", "unknown")

    if not text or not text.strip():
        return []

    # Split text into logical sections before applying character limits
    raw_sections = _split_into_sections(text, extracted)

    # Now chunk sections, merging small ones and splitting large ones
    chunks: List[dict] = []
    current_chunk = ""
    current_sections: List[str] = []

    for section in raw_sections:
        if len(current_chunk) + len(section) + 2 <= max_chars:  # +2 for \n\n join
            current_chunk += ("\n\n" if current_chunk else "") + section
            current_sections.append(section[:60])  # label hint
        else:
            # Save current chunk if non-empty
            if current_chunk.strip():
                chunks.append({
                    "text": current_chunk.strip(),
                    "chunk_index": len(chunks),
                    "total_chunks": -1,  # filled below
                    "source_doc": source_doc,
                    "section_hint": "; ".join(current_sections[:3]),
                    "extraction_method": method,
                    "quality": quality,
                })

            # If single section > max_chars, hard-split it
            if len(section) > max_chars:
                hard_chunks = _hard_split(section, max_chars)
                for hc in hard_chunks:
                    chunks.append({
                        "text": hc.strip(),
                        "chunk_index": len(chunks),
                        "total_chunks": -1,
                        "source_doc": source_doc,
                        "section_hint": section[:60],
                        "extraction_method": method,
                        "quality": quality,
                    })
                current_chunk = ""
                current_sections = []
            else:
                current_chunk = section
                current_sections = [section[:60]]

    # Don't forget the last chunk
    if current_chunk.strip():
        chunks.append({
            "text": current_chunk.strip(),
            "chunk_index": len(chunks),
            "total_chunks": -1,
            "source_doc": source_doc,
            "section_hint": "; ".join(current_sections[:3]),
            "extraction_method": method,
            "quality": quality,
        })

    # Fill in total_chunks now that we know it
    total = len(chunks)
    for c in chunks:
        c["total_chunks"] = total

    return chunks


def _split_into_sections(text: str, extracted: dict) -> List[str]:
    """
    Split text into logical sections based on document type.

    Returns list of section strings. Each section represents a natural
    boundary in the document (sheet, slide, or paragraph).
    """
    method = extracted.get("extraction_method", "")

    # Excel: already sectioned by sheet — split on sheet markers
    if "openpyxl" in method or "pandas" in method:
        sections = [s.strip() for s in text.split("[Sheet:") if s.strip()]
        return ["[Sheet:" + s for s in sections] if sections else [text]

    # PPTX: already sectioned by slide — split on slide markers
    if "pptx" in method:
        sections = [s.strip() for s in text.split("[Slide") if s.strip()]
        return ["[Slide " + s for s in sections] if sections else [text]

    # PDF and DOCX: split on double newlines (paragraph/section boundaries)
    # This preserves paragraph integrity — never breaks mid-paragraph
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    return paragraphs if paragraphs else [text]


def _hard_split(text: str, max_chars: int) -> List[str]:
    """
    Last-resort split for a single section that exceeds max_chars.

    Splits at sentence boundaries where possible (. ! ? or newline),
    falling back to the character limit if no boundary is found within
    500 chars of the limit.
    """
    if len(text) <= max_chars:
        return [text]

    chunks: List[str] = []
    while len(text) > max_chars:
        # Try to find a sentence boundary near the limit
        split_at = max_chars
        for i in range(max_chars, max(0, max_chars - 500), -1):
            if text[i] in ".!?\n":
                split_at = i + 1
                break
        chunks.append(text[:split_at].strip())
        text = text[split_at:].strip()

    if text:
        chunks.append(text)

    return chunks
